"""Extract slide text and speaker notes from modern PowerPoint files."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.ingestion.pdf_parser import PdfDocument, PdfPage


def parse_pptx(path: str | Path) -> PdfDocument:
    """Represent each slide as a numbered page for the shared chunker."""
    source = Path(path)
    presentation = Presentation(str(source))
    pages: list[PdfPage] = []

    for number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                parts.append(f"Speaker notes:\n{notes.strip()}")
        except (AttributeError, ValueError):
            pass
        pages.append(PdfPage(page_number=number, text="\n\n".join(parts)))

    return PdfDocument(
        path=source,
        pages=pages,
        info={"format": "PowerPoint", "slides": len(pages)},
    )
