"""PowerPoint lecture-material extraction tests."""

from pptx import Presentation

from app.ingestion.pptx_parser import parse_pptx


def test_parse_pptx_extracts_slide_text(tmp_path):
    path = tmp_path / "lecture.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Human-computer interaction"
    slide.placeholders[1].text = "Affordances and feedback"
    deck.save(path)

    parsed = parse_pptx(path)

    assert len(parsed.pages) == 1
    assert parsed.pages[0].page_number == 1
    assert "Human-computer interaction" in parsed.pages[0].text
    assert "Affordances and feedback" in parsed.pages[0].text
